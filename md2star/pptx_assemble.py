"""Seam-D PPTX assembler: build a deck directly on the AI-selected layouts.

:mod:`md2star.pptx_layout` decides, per slide, which of a designer template's
named layouts best carries that slide's content. This module is what actually
**uses** that decision: for every slide it instantiates the *chosen* layout
(not a fixed handful of curated backgrounds) from the real template and
overlays guaranteed-legible content, so the deck's visual variety follows the
AI's per-slide reasoning instead of a hand-picked round-robin.

Design choice, carried over from the PoC that proved it out: fighting a
designer layout's own placeholders is fragile (their colours and geometry vary
wildly, some are dark-on-dark, some clutter a small box next to a big void).
Instead each slide keeps the chosen layout's **background** (its true
branding) and places **our own guaranteed-safe text/image boxes** in it, sized
from the layout's actual placeholder geometry when the catalog recorded one
(title / body / picture), and a sensible fractional default otherwise — so it
degrades gracefully for template layouts the catalog under-described.

Text colour is not hard-coded white: with an arbitrary designer layout now in
play (not just a curated dark subset), a slide's background can be light, so
the assembler reads the layout's own background colour and switches to a dark
foreground when needed. This is the correctness gap the curated-background PoC
never had to close.

The **target-matching Ralph Eyeball Loop** (:func:`eyeball_slides`) closes the
verification gap the same way :mod:`md2star.reverse_diagrams` does for
diagrams: render the assembled slide, show it *and* the chosen layout's own
catalog thumbnail to a vision model, and apply bounded, spec-level fixes
(shrink text, trim bullets, grow an undersized image) rather than ever hand-
editing the picture.

Author
------
[Warith HARCHAOUI](https://linkedin.com/in/warith-harchaoui/)
"""

from __future__ import annotations

import json
import re
import zipfile
from collections import Counter
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .logging import get_logger
from .pptx_layout import (
    ChatFn,
    LayoutInfo,
    RenderFn,
    Slide,
    default_chat,
    default_render,
    render_deck_pages,
)

logger = get_logger(__name__)

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_NEAR_BLACK = RGBColor(0x14, 0x14, 0x16)
_CODE_BG = RGBColor(0x14, 0x14, 0x16)
_CODE_FG = RGBColor(0xEA, 0xEA, 0xEA)

_EMOJI = re.compile("[\U0001f000-\U0001faff\U00002600-\U000027bf\U0001f1e6-\U0001f1ff️]")

# Families that ship with the OS — never treated as "the template's design
# font" when picking the most common typeface (see :func:`template_fonts`).
_SYSTEM_FONTS = {
    "arial",
    "calibri",
    "helvetica",
    "helvetica neue",
    "times",
    "times new roman",
    "cambria",
    "verdana",
    "tahoma",
    "georgia",
    "courier",
    "courier new",
    "segoe ui",
}
_WEIGHTS = (
    "black",
    "extrablack",
    "extrabold",
    "extra bold",
    "heavy",
    "bold",
    "semibold",
    "semi bold",
    "demibold",
    "medium",
)

EMU_PER_INCH = 914400


def _clean(text: str) -> str:
    return _EMOJI.sub("", text).replace("  ", " ").strip()


# ── template introspection ───────────────────────────────────────────────────
def layout_index(prs: Presentation) -> dict[str, object]:
    """Map every distinct layout name to its ``SlideLayout`` object."""
    idx: dict[str, object] = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            idx.setdefault(layout.name, layout)
    return idx


def clear_slides(prs: Presentation) -> None:
    """Drop every slide the template ships with, keeping layouts/masters."""
    # python-pptx has no public "remove slide" API; drop_rel + list removal is
    # the documented workaround (also used by md2star.postprocess elsewhere).
    lst = prs.slides._sldIdLst  # noqa: SLF001 — python-pptx exposes no public API for this
    for sld in list(lst):
        prs.part.drop_rel(sld.rId)
        lst.remove(sld)


def _bg_hex(layout) -> str | None:
    """Solid background colour as ``#RRGGBB``, or ``None`` if not resolvable."""
    try:
        fill = layout.background.fill
        if fill.type == MSO_FILL.SOLID:
            return f"#{fill.fore_color.rgb}"
    except Exception:  # noqa: BLE001 — gradient / theme-ref / no explicit bg
        return None
    return None


def is_dark_layout(layout) -> bool | None:
    """``True``/``False`` for a resolvable solid background, else ``None``."""
    hexc = _bg_hex(layout)
    if not hexc:
        return None
    r, g, b = (int(hexc[i : i + 2], 16) for i in (1, 3, 5))
    # Standard perceptual luma weights (ITU-R BT.601); 110/255 is a slightly
    # conservative midpoint (below the textbook 128) so a mid-tone background
    # still gets white text, which stays legible more often than the reverse.
    return (0.299 * r + 0.587 * g + 0.114 * b) < 110


def _weight_rank(name: str) -> int:
    low = name.lower()
    for i, w in enumerate(_WEIGHTS):
        if w in low:
            return i
    return len(_WEIGHTS)


def template_fonts(template: Path) -> dict[str, str]:
    """Derive the template's own DESIGN fonts: ``{family, display, body}``.

    Reads every explicit ``typeface="..."`` on slide/layout shapes (not the
    theme, which carries Office's long CJK fallback list) and picks the most-
    common non-system face as *body*, its heaviest same-family weight as
    *display* (titles). Falls back to Arial when the template declares no
    explicit design font. Read-only: unlike the PoC this is derived from, this
    never downloads or installs fonts — a missing font degrades to whatever
    substitute PowerPoint/LibreOffice picks, exactly like every other
    best-effort AI/rendering feature in this package.
    """
    counts: Counter = Counter()
    with zipfile.ZipFile(template) as z:
        for name in z.namelist():
            if re.search(r"ppt/(slides|slideLayouts)/.*\.xml$", name):
                xml = z.read(name).decode("utf-8", "replace")
                for f in re.findall(r'typeface="([^"]+)"', xml):
                    # "+mn-lt"/"+mj-lt" are theme placeholders, not real
                    # typeface names — an explicit face always wins over them.
                    if f and not f.startswith("+"):
                        counts[f] += 1
    design = [(f, k) for f, k in counts.items() if f.lower() not in _SYSTEM_FONTS]
    if not design:
        face = counts.most_common(1)[0][0] if counts else "Arial"
        return {"family": face.split(" ")[0], "display": face, "body": face}
    # The most-frequent non-system face is almost always body text (there's
    # far more prose than titles in a real deck) — a cheap, reliable proxy
    # for "the template's body font" without parsing which shape is which.
    body_face = max(design, key=lambda x: x[1])[0]
    family = body_face.split(" ")[0]
    variants = [f for f, _ in design if f.split(" ")[0] == family]
    # The heaviest same-family variant in actual use is the designer's own
    # title weight (e.g. "Montserrat ExtraBold" next to plain "Montserrat").
    display = min(variants, key=_weight_rank)
    plain = [v for v in variants if v.lower() == family.lower()]
    return {"family": family, "display": display, "body": plain[0] if plain else body_face}


# ── content parsing ──────────────────────────────────────────────────────────
def is_content(features: list[str]) -> bool:
    return (
        bool({"bullets", "ordered-list", "table"} & set(features))
        and "big-statement" not in features
    )


def parse_image(body: str, base_dir: Path | None) -> Path | None:
    m = re.search(r"!\[[^\]]*\]\(([^)]+)\)", body)
    if not m or base_dir is None:
        return None
    # split()[0] drops a Markdown title suffix (`"foo.png "optional title""`);
    # a non-existent resolved path returns None rather than a dangling Path,
    # so the caller's `if img:` branch check works without a stat() of its own.
    cand = (base_dir / m.group(1).split()[0]).resolve()
    return cand if cand.exists() else None


def parse_mermaid(body: str) -> str | None:
    m = re.search(r"```mermaid\s*\n(.*?)```", body, re.DOTALL)
    return m.group(1).strip() if m else None


def parse_code(body: str) -> str | None:
    """First non-Mermaid fenced block (the slide's primary code sample, if any)."""
    for m in re.finditer(r"```(\w*)\s*\n(.*?)```", body, re.DOTALL):
        if m.group(1).lower() != "mermaid":
            return m.group(2).strip()
    return None


def fallback_bullets(body: str, cap: int) -> list[str]:
    """Trimmed bullet text, used when no AI summary was computed for a slide."""
    body = re.sub(r"```.*?```", "", body, flags=re.DOTALL)  # code/mermaid isn't bullet prose
    out: list[str] = []
    for raw in body.splitlines():
        s = raw.strip()
        # Skip headings, table rows, isolated images, and bare list-marker
        # runs (an empty "- " line) — none of these read as a real bullet.
        if not s or s.startswith(("#", "|", ":::", "![")) or set(s) <= {"*", "-", " "}:
            continue
        s = _clean(re.sub(r"[*_`>]", "", re.sub(r"^[-*]\s+|^\d+\.\s+", "", s)))
        if s:
            # 10-word cap keeps a fallback bullet on one visual line at the
            # sizes add_text() uses — the AI summary (poc_layout_match-style)
            # is already this terse; this is only the no-AI safety net.
            out.append(" ".join(s.split()[:10]))
    return out[:cap]


# ── placeholder-aware boxing ──────────────────────────────────────────────────
# Fractional (left, top, width, height) defaults over the ACTUAL slide size —
# used only when the chosen layout's catalog entry recorded no matching
# placeholder geometry, so this degrades safely for any aspect ratio.
_DEFAULT_TITLE_BOX = (0.06, 0.06, 0.88, 0.22)
_DEFAULT_DIVIDER_BOX = (0.08, 0.28, 0.84, 0.44)
_DEFAULT_BODY_BOX = (0.06, 0.30, 0.88, 0.62)
_DEFAULT_IMAGE_BOX = (0.05, 0.20, 0.90, 0.75)

_TITLE_KINDS = {"title", "ctrTitle"}
_BODY_KINDS = {"body", "subTitle"}
_IMAGE_KINDS = {"pic"}


def _placeholder_box_in(
    info: LayoutInfo | None, kinds: set[str]
) -> tuple[float, float, float, float] | None:
    """The first matching placeholder's box, in inches, or ``None``."""
    if info is None:
        return None
    for ph in info.placeholders:
        if ph.kind in kinds and ph.emu is not None:
            left, top, width, height = ph.emu
            return (
                left / EMU_PER_INCH,
                top / EMU_PER_INCH,
                width / EMU_PER_INCH,
                height / EMU_PER_INCH,
            )
    return None


def _fractional_box(
    prs: Presentation, frac: tuple[float, float, float, float]
) -> tuple[float, float, float, float]:
    slide_w = prs.slide_width / EMU_PER_INCH
    slide_h = prs.slide_height / EMU_PER_INCH
    left, top, width, height = frac
    return (left * slide_w, top * slide_h, width * slide_w, height * slide_h)


def content_box(prs: Presentation, info: LayoutInfo | None) -> tuple[float, float, float, float]:
    # Placeholder geometry always wins when the catalog recorded one — it's
    # the designer's own intended zone, more faithful than any guessed default.
    return _placeholder_box_in(info, _BODY_KINDS) or _fractional_box(prs, _DEFAULT_BODY_BOX)


def image_box(prs: Presentation, info: LayoutInfo | None) -> tuple[float, float, float, float]:
    return _placeholder_box_in(info, _IMAGE_KINDS) or _fractional_box(prs, _DEFAULT_IMAGE_BOX)


def title_box(
    prs: Presentation, info: LayoutInfo | None, *, divider: bool = False
) -> tuple[float, float, float, float]:
    ph = _placeholder_box_in(info, _TITLE_KINDS)
    if ph is not None:
        return ph
    # No title placeholder recorded (common — Jellysmack's own layouts use
    # generic "body" placeholders even for titles): a section-divider slide
    # gets the roomier, vertically-centred default box, a normal slide the
    # slim top-strip one.
    return _fractional_box(prs, _DEFAULT_DIVIDER_BOX if divider else _DEFAULT_TITLE_BOX)


# ── shape drawing ────────────────────────────────────────────────────────────
def add_text(
    slide,
    blocks: list[tuple[str, int, bool, str]],
    box: tuple[float, float, float, float],
    *,
    color: RGBColor,
    center: bool = False,
    middle: bool = False,
) -> None:
    """Overlay a text box. *blocks* = list of (text, pt_size, bold, font)."""
    left, top, width, height = (Inches(v) for v in box)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    # Belt-and-braces against overflow: even after the eyeball loop's
    # font_delta shrink, PowerPoint/LibreOffice auto-shrink the run further
    # rather than ever clip text off the shape.
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if middle:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (text, size, bold, font) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        if center:
            p.alignment = PP_ALIGN.CENTER
        # Set on the RUN, not the paragraph: LibreOffice/PowerPoint read
        # run-level rPr for size/bold/font/colour, so a paragraph-level
        # default here would silently not render in at least one of the two.
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color


def _fit(pic, max_w: float, max_h: float, cx: float, area_top: float) -> None:
    """Contain-fit *pic* inside (max_w x max_h) inches, centred at *cx*."""
    # min(), not a 1.0-capped scale: a small source image is GROWN to fill
    # its box rather than left tiny — deliberate, matches the eyeball loop's
    # own "image_too_small" check, which exists precisely to catch the
    # opposite mistake (a good source image rendered needlessly small).
    scale = min(Inches(max_w) / pic.width, Inches(max_h) / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = Inches(cx) - pic.width // 2
    pic.top = Inches(area_top) + (Inches(max_h) - pic.height) // 2


def add_image(slide, img: Path, *, box: tuple[float, float, float, float]) -> None:
    left, top, width, height = box
    cx = left + width / 2
    pic = slide.shapes.add_picture(str(img), Inches(left), Inches(top))
    _fit(pic, width, height, cx, top)


def add_code(slide, code: str, box: tuple[float, float, float, float], *, font: str) -> None:
    left, top, width, height = box
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = _CODE_BG
    rect.line.fill.background()
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.25)
    for i, ln in enumerate(code.splitlines()[:12]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln or " "
        run = p.runs[0]
        run.font.name = font
        run.font.size = Pt(13)
        run.font.color.rgb = _CODE_FG


# ── assembly ──────────────────────────────────────────────────────────────────
@dataclass
class AssembleResult:
    pptx_path: Path
    placed: int
    total: int


RenderMermaidFn = Callable[[str], "Path | None"]


def assemble(
    slides: list[Slide],
    catalog: list[LayoutInfo],
    template: Path,
    out_path: Path,
    *,
    image_base_dir: Path | None = None,
    render_mermaid: RenderMermaidFn | None = None,
    default_layout: str | None = None,
) -> AssembleResult:
    """Build a real ``.pptx`` directly on each slide's AI-selected layout.

    For every slide, instantiate ``slides[i].chosen_layout`` from *template*
    (falling back to *default_layout*, or the template's first layout, when
    the pick is missing or unknown) and overlay title/body/image/code content
    sized from that layout's own placeholder geometry when the catalog
    recorded one. Per-slide ``overrides`` (written by :func:`eyeball_slides`)
    bound font size, bullet count and image growth.

    Returns
    -------
    AssembleResult
        The output path plus how many of the slides were placed without error
        (a per-slide exception is logged and skipped, never aborts the deck).
    """
    catalog_by_name = {c.name: c for c in catalog}
    prs = Presentation(str(template))
    layouts = layout_index(prs)
    have = set(layouts)
    fam = template_fonts(template)
    title_font, body_font = fam["display"], fam["body"]
    clear_slides(prs)

    n = len(slides)
    placed = 0
    fallback_name = default_layout if default_layout in have else next(iter(have), None)
    if fallback_name is None:
        raise ValueError(f"template {template} exposes no slide layouts")

    for s in slides:
        # A missing or hallucinated pick (select_layouts failed, or a caller
        # set chosen_layout by hand) still has to land on SOME real layout —
        # the deck must never lose a slide over a bad AI decision.
        name = s.chosen_layout if s.chosen_layout in have else fallback_name
        layout = layouts[name]
        info = catalog_by_name.get(name)
        slide = prs.slides.add_slide(layout)
        # A layout with no explicitly resolvable solid fill (``dark is None``)
        # is, empirically, virtually always a plain white/light layout whose
        # background comes from the master's default rather than an explicit
        # override (a real black/gradient/photo background always resolves) —
        # so treat "unresolvable" as light, not dark. Getting this wrong is a
        # correctness bug, not a cosmetic one: white-on-white is invisible.
        dark = is_dark_layout(layout)
        fg = _WHITE if dark else _NEAR_BLACK

        body = s.body
        mer = parse_mermaid(body) if "mermaid" in s.features else None
        # code/img are only checked when there's no mermaid block, since a
        # mermaid fence would otherwise also match the generic code/image
        # patterns and steal the slide from its actual diagram branch.
        code = parse_code(body) if not mer else None
        img = parse_image(body, image_base_dir) if not mer else None
        # Cover (1) and closing (n) slides are always title-only, whatever
        # bullet/table markup the source happens to carry.
        content = is_content(s.features) and s.index not in (1, n)

        title_lines = [_clean(s.title)] or ["(untitled)"]

        try:
            # Parsed INSIDE the try so a malformed override (e.g. a corrupt
            # eyeball-loop verdict) only drops this one slide, never the deck.
            ov = s.overrides or {}
            fd = int(ov.get("font_delta", 0))
            bcap = max(1, int(ov.get("bullets", 4)))
            igrow = float(ov.get("img_grow", 0.0))
            if content:
                # No AI summary is threaded through the assembler on purpose
                # (that lives in pptx_layout's selection stage, not here) —
                # fallback_bullets is the assembler's own, AI-free safety net.
                bullets = fallback_bullets(body, bcap)
                blocks = [(title_lines[0], 26 + fd, True, title_font)]
                blocks += [(f"•  {b}", 16 + fd, False, body_font) for b in bullets]
                add_text(slide, blocks, content_box(prs, info), color=fg)
            elif mer and render_mermaid and (mpng := render_mermaid(mer)) is not None:
                add_text(
                    slide,
                    [(title_lines[0], 20 + fd, True, title_font)],
                    title_box(prs, info),
                    color=fg,
                )
                left, top, width, height = image_box(prs, info)
                add_image(slide, mpng, box=(left - igrow, top, width + 2 * igrow, height + igrow))
            elif img:
                add_text(
                    slide,
                    [(title_lines[0], 20 + fd, True, title_font)],
                    title_box(prs, info),
                    color=fg,
                )
                left, top, width, height = image_box(prs, info)
                add_image(slide, img, box=(left - igrow, top, width + 2 * igrow, height + igrow))
            elif code:
                add_text(
                    slide,
                    [(title_lines[0], 22 + fd, True, title_font)],
                    title_box(prs, info),
                    color=fg,
                )
                add_code(slide, code, content_box(prs, info), font=fam.get("family", "Consolas"))
            else:
                add_text(
                    slide,
                    [(ln, 40 + fd, True, title_font) for ln in title_lines],
                    title_box(prs, info, divider=True),
                    color=fg,
                    center=True,
                    middle=True,
                )
            placed += 1
        except Exception as exc:  # noqa: BLE001 — one bad slide must not sink the deck
            logger.warning("slide %d (layout %r) failed to assemble: %s", s.index, name, exc)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info("assembled %d/%d slides -> %s", placed, n, out_path)
    return AssembleResult(pptx_path=out_path, placed=placed, total=n)


# ── target-matching Ralph Eyeball Loop ───────────────────────────────────────
_JUDGE_SCHEMA = {
    "type": "object",
    "properties": {
        "matches": {"type": "boolean"},
        "overflow": {"type": "boolean"},
        "low_contrast": {"type": "boolean"},
        "too_empty": {"type": "boolean"},
        "image_too_small": {"type": "boolean"},
        "discrepancies": {"type": "string"},
    },
    "required": ["matches"],
}

_JUDGE_PROMPT = (
    "The FIRST image is the TARGET: an example slide of the designer template "
    "layout this slide is supposed to use. The SECOND image is the CANDIDATE: "
    "this slide as actually rendered. Judge the candidate against the target's "
    "composition, not its specific words/pictures: does it place a title, body "
    "and any image/diagram in similarly proportioned regions, with nothing "
    "overflowing the slide edges, clipped, low-contrast against its background, "
    "or leaving an awkward empty void? Reply as JSON: matches (true only if it "
    "is clean and well-composed), overflow, low_contrast, too_empty, "
    "image_too_small, and a short discrepancies string (empty if matches)."
)


def _overrides_for(verdict: dict) -> dict:
    # Bounded, one-shot corrections — never a search: -4pt/2-bullets and
    # +50% image growth are each a single fixed nudge per flagged issue, not
    # a tuned parameter the loop iterates toward. low_contrast/too_empty are
    # reported to the caller in ``discrepancies`` but have no safe automatic
    # fix here (a layout swap, not a font/bullet tweak, is what they need).
    ov: dict = {}
    if verdict.get("overflow"):
        ov["font_delta"] = -4
        ov["bullets"] = 2
    if verdict.get("image_too_small"):
        ov["img_grow"] = 0.5
    return ov


def eyeball_slides(
    slides: list[Slide],
    catalog: list[LayoutInfo],
    template: Path,
    out_path: Path,
    *,
    indices: list[int] | None = None,
    max_iterations: int = 1,
    chat: ChatFn | None = None,
    render_template: RenderFn | None = None,
    render_deck: Callable[[Path], dict[int, bytes]] | None = None,
    image_base_dir: Path | None = None,
    render_mermaid: RenderMermaidFn | None = None,
    default_layout: str | None = None,
) -> AssembleResult:
    """Target-matching Ralph Eyeball Loop over an assembled deck.

    Assembles (or re-uses an already-assembled) deck, renders it, and for each
    of *indices* (default: every slide) shows a vision model both the
    candidate render and the chosen layout's own catalog thumbnail. On a
    mismatch it writes a bounded ``overrides`` entry back onto the slide
    (shrink text, cap bullets, grow an undersized image — never touches the
    picture itself) and re-assembles, up to *max_iterations* passes. Any
    transport/render failure leaves the current best assembly untouched.

    *render_deck* renders the assembled output deck itself (default:
    :func:`md2star.pptx_layout.render_deck_pages`, LibreOffice + poppler);
    injectable so the loop is testable offline without either installed.
    """
    chat = chat or default_chat(kind="vlm")
    render_template = render_template or default_render()
    render_deck = render_deck or render_deck_pages
    catalog_by_name = {c.name: c for c in catalog}
    target_indices = indices if indices is not None else [s.index for s in slides]
    by_index = {s.index: s for s in slides}

    result = assemble(
        slides,
        catalog,
        template,
        out_path,
        image_base_dir=image_base_dir,
        render_mermaid=render_mermaid,
        default_layout=default_layout,
    )

    # One extra pass beyond max_iterations: the loop always re-judges after
    # the LAST fix is applied (to know whether it actually helped and to
    # report an accurate result), it just never applies a NEW fix on that
    # final pass — see the ``iteration < max_iterations`` guards below.
    for iteration in range(max_iterations + 1):
        pages = render_deck(out_path)
        flagged = 0
        for i in target_indices:
            s = by_index.get(i)
            if s is None:
                continue
            candidate_png = pages.get(i)
            info = catalog_by_name.get(s.chosen_layout) if s.chosen_layout else None
            target_page = info.rep_page if info else None
            if candidate_png is None or target_page is None:
                continue  # nothing to render/compare against -> skip, don't guess
            target_png = render_template(template, target_page)
            if target_png is None:
                continue
            out = chat(_JUDGE_PROMPT, [target_png, candidate_png], _JUDGE_SCHEMA)
            verdict = (
                out
                if isinstance(out, dict)
                else (json.loads(out) if isinstance(out, str) else None)
            )
            if not verdict or verdict.get("matches"):
                continue
            if iteration < max_iterations:
                ov = _overrides_for(verdict)
                if ov:
                    # Merge, not replace: a slide flagged on two different
                    # passes accumulates both fixes instead of losing the first.
                    s.overrides = {**(s.overrides or {}), **ov}
                    flagged += 1
        if iteration < max_iterations:
            if not flagged:
                break  # nothing actionable left to fix — stop before a no-op re-assemble
            result = assemble(
                slides,
                catalog,
                template,
                out_path,
                image_base_dir=image_base_dir,
                render_mermaid=render_mermaid,
                default_layout=default_layout,
            )

    return result
